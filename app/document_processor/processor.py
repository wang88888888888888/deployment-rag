from typing import List, Dict, Any
import PyPDF2
from pptx import Presentation
from pathlib import Path
from langchain.text_splitter import RecursiveCharacterTextSplitter

class DocumentProcessor:
    def __init__(self, chunk_size: int = 500, chunk_overlap: int = 50):
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            length_function=len,
            separators=["\n\n", "\n", " ", ""]
        )
        
    def process_pdf(self, file_path: str) -> List[Dict[str, Any]]:
        """Process PDF file and return chunks with metadata."""
        chunks = []
        extracted_any = False
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page_num, page in enumerate(pdf_reader.pages):
                text = page.extract_text()
                if not text or not text.strip():
                    print(f"[WARN] No text extracted from page {page_num + 1} of {file_path}")
                    continue
                extracted_any = True
                print(f"[INFO] Extracted text from page {page_num + 1} of {file_path} (first 100 chars): {text[:100]}")
                page_chunks = self.text_splitter.split_text(text)
                
                for chunk in page_chunks:
                    chunks.append({
                        "content": chunk,
                        "metadata": {
                            "source": Path(file_path).name,
                            "page": page_num + 1,
                            "type": "pdf"
                        }
                    })
        if not extracted_any:
            print(f"[WARN] No text extracted from any page in PDF: {file_path}")
        return chunks
    
    def process_pptx(self, file_path: str) -> List[Dict[str, Any]]:
        """Process PowerPoint file and return chunks with metadata."""
        chunks = []
        prs = Presentation(file_path)
        
        for slide_num, slide in enumerate(prs.slides):
            text_content = []
            # Extract slide title
            if slide.shapes and hasattr(slide.shapes[0], "text"):
                title = slide.shapes[0].text
                if title:
                    text_content.append(f"Title: {title}")
            # Extract all text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_content.append(shape.text)
                # Extract alt text/caption if available
                if hasattr(shape, "alternative_text") and shape.alternative_text:
                    text_content.append(f"Alt text: {shape.alternative_text}")
                if hasattr(shape, "caption") and shape.caption:
                    text_content.append(f"Caption: {shape.caption}")
            # Extract notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text
                if notes:
                    text_content.append(f"Notes: {notes}")
            slide_text = "\n".join(text_content)
            if not slide_text.strip():
                print(f"[WARN] No text extracted from slide {slide_num + 1} of {file_path}")
                continue
            print(f"[INFO] Extracted text from slide {slide_num + 1} of {file_path} (first 100 chars): {slide_text[:100]}")
            slide_chunks = self.text_splitter.split_text(slide_text)
            
            for chunk in slide_chunks:
                chunks.append({
                    "content": chunk,
                    "metadata": {
                        "source": Path(file_path).name,
                        "slide": slide_num + 1,
                        "type": "pptx"
                    }
                })
        return chunks
    
    def process_document(self, file_path: str) -> List[Dict[str, Any]]:
        """Process document based on file extension."""
        file_path = Path(file_path)
        if file_path.suffix.lower() == '.pdf':
            return self.process_pdf(str(file_path))
        elif file_path.suffix.lower() in ['.pptx', '.ppt']:
            return self.process_pptx(str(file_path))
        else:
            raise ValueError(f"Unsupported file type: {file_path.suffix}") 